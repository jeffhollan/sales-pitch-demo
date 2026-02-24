"""Document generation tools — Word prep doc and branded PowerPoint.

These functions are structured as Copilot SDK function tools with annotated
parameters so they can be registered via ``GitHubCopilotAgent(tools=[...])``.
"""

from __future__ import annotations

import json
import zipfile
from datetime import date
from io import BytesIO
from pathlib import Path
from typing import Annotated, Any

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from pydantic import Field

from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches, Pt as PptxPt

from src.config import MOCK_DATA_DIR, PPTX_TEMPLATE, ensure_output_dir


# ── Brand helpers ──────────────────────────────────────────────────────

def _load_brand(customer_name: str) -> dict[str, Any]:
    brands_file = MOCK_DATA_DIR / "brands.json"
    with open(brands_file) as f:
        brands = json.load(f)
    key = customer_name.lower().replace("the ", "").replace(" company", "").replace(" ", "-").strip()
    for k, v in brands.items():
        if k in key or key in k:
            return v
    return {
        "display_name": customer_name,
        "short_name": customer_name,
        "primary_color": "#0078D4",
        "secondary_color": "#000000",
        "accent_color": "#FFFFFF",
        "industry": "Unknown",
    }


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


# ── Microsoft brand colors ────────────────────────────────────────────

MS_BLUE = RGBColor(0, 120, 212)       # #0078D4
MS_DARK = RGBColor(36, 36, 36)        # #242424
MS_GRAY = RGBColor(96, 96, 96)        # #606060
MS_LIGHT_GRAY = RGBColor(200, 200, 200)


# ── Template loader ───────────────────────────────────────────────────

def _load_template(template_path: Path | None = None) -> Presentation:
    """Load the Microsoft brand .potx template as a Presentation with no slides.

    python-pptx rejects .potx content types, so we patch the ZIP in memory:
    1. Change the content-type from *template* to *presentation*.
    2. Strip all existing slide parts and their relationships so the
       Presentation loads with layouts/masters only — no sample slides.
    """
    import re

    path = template_path or PPTX_TEMPLATE
    with open(path, "rb") as f:
        raw = f.read()

    src = BytesIO(raw)
    dst = BytesIO()
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(dst, "w") as zout:
        for item in zin.infolist():
            # Skip all slide parts and their rels (but keep slideMasters / slideLayouts)
            if re.match(r"ppt/slides/", item.filename):
                continue
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                # Fix content type for .potx → .pptx
                data = data.replace(
                    b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                )
                # Remove content-type entries for slide parts
                data = re.sub(
                    rb'<Override[^>]*PartName="/ppt/slides/[^"]*"[^>]*/>\s*',
                    b"",
                    data,
                )
            if item.filename == "ppt/presentation.xml":
                # Remove <p:sldIdLst>...</p:sldIdLst> so no slides are referenced
                data = re.sub(
                    rb"<p:sldIdLst>.*?</p:sldIdLst>",
                    b"<p:sldIdLst/>",
                    data,
                    flags=re.DOTALL,
                )
            if item.filename == "ppt/_rels/presentation.xml.rels":
                # Remove relationship entries that point to slides/
                data = re.sub(
                    rb'<Relationship[^>]*Target="slides/[^"]*"[^>]*/>\s*',
                    b"",
                    data,
                )
            zout.writestr(item, data)
    dst.seek(0)
    return Presentation(dst)


# ═══════════════════════════════════════════════════════════════════════
#  WORD PREP DOC
# ═══════════════════════════════════════════════════════════════════════

def generate_prep_doc(
    customer_name: Annotated[str, Field(description="Customer company name")],
    work_iq: Annotated[dict[str, Any] | str, Field(description="Work IQ data (dict or JSON string)")],
    fabric_iq: Annotated[dict[str, Any] | str, Field(description="Fabric IQ data (dict or JSON string)")],
    foundry_iq: Annotated[dict[str, Any] | str, Field(description="Foundry IQ data (dict or JSON string)")],
) -> str:
    """Generate a Word meeting prep document for a customer meeting.

    Returns the output file path.
    """
    # Accept JSON strings (for function-tool invocation) or dicts
    if isinstance(work_iq, str):
        work_iq = json.loads(work_iq)
    if isinstance(fabric_iq, str):
        fabric_iq = json.loads(fabric_iq)
    if isinstance(foundry_iq, str):
        foundry_iq = json.loads(foundry_iq)

    out_dir = ensure_output_dir()
    brand = _load_brand(customer_name)
    doc = Document()

    # -- Page style --
    style = doc.styles["Normal"]
    style.font.name = "Segoe UI"
    style.font.size = Pt(10)
    style.font.color.rgb = MS_DARK

    for level in range(1, 4):
        hs = doc.styles[f"Heading {level}"]
        hs.font.name = "Segoe UI Semibold"
        hs.font.color.rgb = MS_BLUE

    # -- Title --
    title = doc.add_heading(f"Meeting Prep — {brand['display_name']}", level=0)
    title.runs[0].font.color.rgb = MS_BLUE
    title.runs[0].font.size = Pt(22)
    doc.add_paragraph(f"Prepared {date.today().strftime('%B %d, %Y')}  |  Confidential — Microsoft Internal")

    doc.add_paragraph("")  # spacer

    # ── Section 1: Relationship Context ──
    doc.add_heading("Relationship Context", level=1)

    contact = work_iq.get("primary_contact") or {}
    if contact.get("name"):
        doc.add_paragraph(
            f"Primary Contact: {contact['name']}, {contact.get('title', 'N/A')}"
        )

    team = work_iq.get("account_team") or {}
    if team:
        doc.add_paragraph(
            f"Account Team: AE {team.get('account_executive', 'N/A')} · "
            f"Tech Lead {team.get('technical_lead', 'N/A')} · "
            f"CSM {team.get('customer_success', 'N/A')}"
        )

    if "relationship_summary" in work_iq:
        doc.add_paragraph(work_iq["relationship_summary"])

    # Recent communications
    if "recent_emails" in work_iq:
        doc.add_heading("Recent Communications", level=2)
        for email in work_iq["recent_emails"][:3]:
            p = doc.add_paragraph()
            run = p.add_run(f"[{email['date']}] {email['subject']}")
            run.bold = True
            run.font.size = Pt(9)
            p.add_run(f"\n{email['snippet']}")
            p.runs[-1].font.size = Pt(9)
            p.runs[-1].font.color.rgb = MS_GRAY

    # ── Section 2: Business Health ──
    doc.add_heading("Business Health", level=1)

    if "financial_summary" in fabric_iq:
        fin = fabric_iq["financial_summary"]
        doc.add_paragraph(
            f"Current Annual Spend: ${fin['current_annual_spend']:,.0f}  |  "
            f"Proposed: ${fin['proposed_annual_spend']:,.0f} (+{fin['growth_potential_pct']}%)  |  "
            f"Renewal Risk: {fin['renewal_risk']}"
        )

    if "usage_trends" in fabric_iq:
        doc.add_heading("Key Usage Metrics", level=2)
        trends = fabric_iq["usage_trends"]
        if "m365_copilot" in trends:
            cp = trends["m365_copilot"]
            doc.add_paragraph(
                f"• M365 Copilot: {cp['active_users_pct']}% active users, "
                f"{cp['monthly_actions_per_user']} actions/user/mo, "
                f"{cp['satisfaction_score']}/5 satisfaction ({cp['trend']})"
            )
        if "fabric" in trends:
            fb = trends["fabric"]
            doc.add_paragraph(
                f"• Fabric: {fb['monthly_queries']:,} queries/mo, "
                f"{fb['capacity_utilization_pct']}% capacity utilization, "
                f"P95 latency {fb['latency_p95_ms']}ms ({fb['latency_trend']})"
            )

    if "support_tickets" in fabric_iq:
        open_tickets = [t for t in fabric_iq["support_tickets"] if t["status"] == "Open"]
        if open_tickets:
            doc.add_heading("Open Support Issues", level=2)
            for t in open_tickets:
                doc.add_paragraph(
                    f"⚠ [{t['severity']}] {t['title']} — opened {t['opened']}\n"
                    f"  {t['description']}"
                )

    # ── Section 3: Recommended Topics ──
    doc.add_heading("Recommended Discussion Topics", level=1)

    if "expansion_opportunities" in fabric_iq:
        for opp in fabric_iq["expansion_opportunities"]:
            doc.add_paragraph(
                f"• {opp['product']} — {opp['stage']} (Confidence: {opp['confidence']})\n"
                f"  Incremental value: ${opp['incremental_value']:,.0f}. {opp.get('notes', '')}",
            )

    # ── Section 4: Relevant Materials ──
    doc.add_heading("Relevant Sales Plays & Materials", level=1)

    if "sales_plays" in foundry_iq:
        for play in foundry_iq["sales_plays"]:
            doc.add_heading(play["play_name"], level=2)
            doc.add_paragraph(play["summary"])
            doc.add_paragraph("Key talking points:")
            for point in play["key_talking_points"]:
                doc.add_paragraph(f"  • {point}")
            if play.get("customer_references"):
                doc.add_paragraph("Customer references:")
                for ref in play["customer_references"]:
                    doc.add_paragraph(f"  → {ref['company']}: {ref['summary']}")

    # -- Save --
    filename = f"meeting_prep_{brand['short_name'].lower().replace(' ', '_')}_{date.today().isoformat()}.docx"
    path = out_dir / filename
    doc.save(str(path))
    return str(path)


# ═══════════════════════════════════════════════════════════════════════
#  POWERPOINT PRESENTATION — Template-based
# ═══════════════════════════════════════════════════════════════════════

# Template layout indices (Microsoft_Brand_Template_May2023.potx)
_LAYOUT_TITLE = 0           # 1_Title_Gradient_Warm Gray
_LAYOUT_BLANK_HEAD = 29     # Blank_with Head (title + custom shapes)
_LAYOUT_3COL_SUBHEADS = 23  # 3-column_Text_with Subheads
_LAYOUT_1COL_TEXT = 19       # 1_1-column_Text


def _set_placeholder_text(slide, idx: int, text: str, *,
                          font_size: int | None = None,
                          bold: bool | None = None,
                          color: str | None = None) -> None:
    """Set text on a template placeholder by its index."""
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if font_size is not None:
        p.font.size = PptxPt(font_size)
    if bold is not None:
        p.font.bold = bold
    if color is not None:
        r, g, b = _hex_to_rgb(color)
        p.font.color.rgb = PptxRGB(r, g, b)


def _add_text_box(slide, left, top, width, height, text, font_size=14,
                  bold=False, color=None, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        PptxInches(left), PptxInches(top), PptxInches(width), PptxInches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PptxPt(font_size)
    p.font.bold = bold
    p.alignment = alignment
    if color:
        r, g, b = _hex_to_rgb(color) if isinstance(color, str) else color
        p.font.color.rgb = PptxRGB(r, g, b)
    return tf


def _add_accent_bar(slide, left, top, width, height, color_hex):
    """Add a colored rectangle accent bar."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        PptxInches(left), PptxInches(top), PptxInches(width), PptxInches(height)
    )
    r, g, b = _hex_to_rgb(color_hex)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptxRGB(r, g, b)
    shape.line.fill.background()
    return shape


def _add_stat_card(slide, left, top, width, height, stat_text, label_text,
                   color_hex):
    """Add a colored stat card with a big number and label overlaid."""
    shape = slide.shapes.add_shape(
        1, PptxInches(left), PptxInches(top),
        PptxInches(width), PptxInches(height),
    )
    r, g, b = _hex_to_rgb(color_hex)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptxRGB(r, g, b)
    shape.line.fill.background()

    _add_text_box(slide, left + 0.15, top + 0.2, width - 0.3, 0.65,
                  stat_text, font_size=44, bold=True, color="#FFFFFF",
                  alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, left + 0.15, top + 0.85, width - 0.3, 0.55,
                  label_text, font_size=12, bold=False, color="#FFFFFF",
                  alignment=PP_ALIGN.CENTER)


def _add_info_card(slide, left, top, width, height, text, *, font_size=11,
                   bg_color=None):
    """Add an info card with body text. Default light-gray background."""
    shape = slide.shapes.add_shape(
        1, PptxInches(left), PptxInches(top),
        PptxInches(width), PptxInches(height),
    )
    r, g, b = bg_color or (235, 235, 235)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptxRGB(r, g, b)
    shape.line.fill.background()

    _add_text_box(slide, left + 0.15, top + 0.15, width - 0.3, height - 0.3,
                  text, font_size=font_size, color="#333333")


# ── Data extraction helpers (internal → customer-facing reframing) ────

def _anonymize_reference(company, summary):
    """Replace company name with an anonymous industry description."""
    labels = {
        "PepsiCo": "A leading CPG company",
        "Nestlé": "Another global CPG leader",
        "AB InBev": "A global brewer",
        "Danone": "A global food & beverage company",
        "Procter & Gamble": "A leading consumer goods company",
    }
    label = labels.get(company, "A leading enterprise")
    anon_summary = summary.replace(company, label) if company in summary else summary
    return f"{label}: {anon_summary}"


def _get_partnership_stats(work_iq, fabric_iq):
    """Extract partnership stats for the 'at a glance' slide."""
    import re as _re
    summary = work_iq.get("relationship_summary", "")
    tenure_match = _re.search(r"(\d+)-year", summary)
    tenure = f"{tenure_match.group(1)} Years" if tenure_match else "Partners"

    contract = fabric_iq.get("contract", {})
    products = contract.get("products", [])

    m365_users = "N/A"
    copilot_users = "N/A"
    for p in products:
        name = p.get("name", "")
        if "365 E5" in name or "365 E3" in name:
            m365_users = f"{p.get('seats', 0):,}"
        if "Copilot" in name:
            copilot_users = f"{p.get('seats', 0):,}"

    return {
        "tenure": tenure,
        "m365_users": m365_users,
        "product_count": str(len(products)),
        "copilot_users": copilot_users,
        "summary": (
            f"A comprehensive technology partnership spanning "
            f"{len(products)} products \u2014 from productivity and "
            f"collaboration to advanced analytics and AI."
        ),
    }


def _get_copilot_highlights(fabric_iq, foundry_iq):
    """Extract Copilot stats and anonymized peer references."""
    cp = fabric_iq.get("usage_trends", {}).get("m365_copilot", {})

    stats = {
        "adoption_pct": f"{cp.get('active_users_pct', 'N/A')}%",
        "actions_per_user": str(cp.get("monthly_actions_per_user", "N/A")),
        "satisfaction": f"{cp.get('satisfaction_score', 'N/A')} / 5.0",
    }

    peers = []
    plays = foundry_iq.get("sales_plays", [])
    if plays:
        refs = plays[0].get("customer_references", [])
        for ref in refs[:2]:
            peers.append(
                _anonymize_reference(ref.get("company", ""), ref.get("summary", ""))
            )
    stats["peers"] = peers
    return stats


def _get_fabric_highlights(fabric_iq, foundry_iq):
    """Extract Fabric stats and optimization commitment."""
    fb = fabric_iq.get("usage_trends", {}).get("fabric", {})

    commitment = ""
    for play in foundry_iq.get("sales_plays", []):
        if "Fabric" in play.get("play_name", ""):
            for point in play.get("key_talking_points", []):
                if "optimization" in point.lower() or "Delta Lake" in point:
                    commitment = point
                    break
            break
    if not commitment:
        commitment = (
            "Implementing performance optimizations for your analytics "
            "pipeline \u2014 creating additional headroom for future growth."
        )

    return {
        "monthly_queries": f"{fb.get('monthly_queries', 0):,}",
        "workloads": fb.get("top_workloads", []),
        "utilization_pct": fb.get("capacity_utilization_pct", 0),
        "commitment": commitment,
    }


def _get_opportunities(fabric_iq, foundry_iq):
    """Build customer-value-framed opportunity columns."""
    import re as _re
    cp = fabric_iq.get("usage_trends", {}).get("m365_copilot", {})
    products = fabric_iq.get("contract", {}).get("products", [])
    plays = foundry_iq.get("sales_plays", [])

    m365_seats = 0
    for p in products:
        if "365 E5" in p.get("name", "") or "365 E3" in p.get("name", ""):
            m365_seats = p.get("seats", 0)

    savings = "$3.2M"
    for play in plays:
        for pt in play.get("key_talking_points", []):
            if "productivity savings" in pt.lower():
                m = _re.search(r"\$[\d.]+M", pt)
                if m:
                    savings = m.group(0)
                break

    adoption = cp.get("active_users_pct", "N/A")

    return [
        {
            "subhead": "AI Copilot for All",
            "body": (
                f"Extend AI-powered productivity to all {m365_seats:,} users. "
                f"Organizations see 15\u201320% less time on routine documents. "
                f"{savings} projected annual savings.\n\n"
                f"Your {adoption}% adoption puts you in the top decile "
                f"of CPG companies."
            ),
        },
        {
            "subhead": "Intelligent Data Platform",
            "body": (
                "Scale your analytics platform for peak performance with "
                "burst capacity and Real-Time Intelligence.\n\n"
                "A leading consumer goods company consolidated 12 analytics "
                "platforms, achieving 40% cost reduction."
            ),
        },
        {
            "subhead": "AI for Operations",
            "body": (
                "AI-powered quality inspection, demand forecasting, and "
                "supply chain optimization for bottling operations.\n\n"
                "A global brewer reduced quality incidents by 35% with "
                "AI-powered visual inspection."
            ),
        },
    ]


def _get_next_steps():
    """Build collaborative next-step action items."""
    return [
        "1. Copilot Expansion Business Case\n"
        "   ROI analysis building on your successful pilot results",

        "2. Data Platform Performance Optimization\n"
        "   Our engineering team is implementing improvements \u2014 "
        "we\u2019ll share results and a scaling plan",

        "3. AI Strategy Workshop\n"
        "   Half-day AI briefing for your operations team \u2014 "
        "exploring quality inspection and demand forecasting",

        "4. Follow-Up\n"
        "   Summary with action items within 48 hours",
    ]


def generate_presentation(
    customer_name: Annotated[str, Field(description="Customer company name")],
    work_iq: Annotated[dict[str, Any] | str, Field(description="Work IQ data as JSON")],
    fabric_iq: Annotated[dict[str, Any] | str, Field(description="Fabric IQ data as JSON")],
    foundry_iq: Annotated[dict[str, Any] | str, Field(description="Foundry IQ data as JSON")],
) -> str:
    """Generate a customer-facing branded PowerPoint presentation.

    Creates a 6-slide deck designed to present TO the customer, framing
    Microsoft capabilities and data as customer value.  All internal data
    (revenue targets, team names, ticket IDs, confidence scores) is omitted
    or reframed.

    Returns the output file path.
    """
    # Accept JSON strings (for function-tool invocation) or dicts
    if isinstance(work_iq, str):
        work_iq = json.loads(work_iq)
    if isinstance(fabric_iq, str):
        fabric_iq = json.loads(fabric_iq)
    if isinstance(foundry_iq, str):
        foundry_iq = json.loads(foundry_iq)

    out_dir = ensure_output_dir()
    brand = _load_brand(customer_name)
    primary = brand["primary_color"]
    ms_blue = "#0078D4"

    contact = work_iq.get("primary_contact", {})
    contact_name = contact.get("name", "")
    contact_title = contact.get("title", "")

    prs = _load_template()

    # ── Slide 1: Title (Layout 0 — 1_Title_Gradient_Warm Gray) ───
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE])
    _add_accent_bar(slide, 0, 0, 13.333, 0.04, primary)

    _set_placeholder_text(
        slide, 0,
        f"Powering {brand['display_name']}\u2019s Digital Transformation",
        font_size=36, bold=True,
    )
    _set_placeholder_text(
        slide, 12,
        f"Microsoft + {brand['display_name']} Partnership Review",
        font_size=20,
    )
    _set_placeholder_text(
        slide, 13,
        date.today().strftime("%B %Y"),
        font_size=14, color="#606060",
    )
    _set_placeholder_text(
        slide, 14,
        f"Prepared for {contact_name}, {contact_title}" if contact_name else "",
        font_size=12, color="#606060",
    )

    # ── Slide 2: Partnership at a Glance ─────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BLANK_HEAD])
    _add_accent_bar(slide, 0, 0, 13.333, 0.04, primary)
    _set_placeholder_text(slide, 0, "Our Partnership at a Glance",
                          font_size=28, bold=True)

    pstats = _get_partnership_stats(work_iq, fabric_iq)

    card_w, card_h = 2.8, 1.5
    gap = 0.3
    total_w = 4 * card_w + 3 * gap
    start_left = (13.333 - total_w) / 2
    card_top = 2.0

    cards = [
        (pstats["tenure"], "Strategic Partnership", ms_blue),
        (pstats["m365_users"], "Users on Microsoft 365", ms_blue),
        (pstats["product_count"], "Microsoft Products", primary),
        (pstats["copilot_users"], "Users with AI Copilot", ms_blue),
    ]
    for i, (stat, label, color) in enumerate(cards):
        left = start_left + i * (card_w + gap)
        _add_stat_card(slide, left, card_top, card_w, card_h, stat, label, color)

    _add_text_box(
        slide, 1.0, card_top + card_h + 0.4, 11.333, 0.6,
        pstats["summary"], font_size=13, color="#606060",
        alignment=PP_ALIGN.CENTER,
    )

    # ── Slide 3: Your Success with AI ────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BLANK_HEAD])
    _add_accent_bar(slide, 0, 0, 13.333, 0.04, primary)
    _set_placeholder_text(slide, 0, "Your Success with AI",
                          font_size=28, bold=True)

    copilot = _get_copilot_highlights(fabric_iq, foundry_iq)

    # Left half — 3 stat callouts stacked vertically
    callouts = [
        (copilot["adoption_pct"], 44, "Copilot adoption across your organization"),
        (copilot["actions_per_user"], 36, "AI-assisted actions per user, per month"),
        (copilot["satisfaction"], 36, "User satisfaction score"),
    ]
    y = 1.8
    for stat_val, fsize, desc in callouts:
        _add_text_box(slide, 0.8, y, 2.0, 0.55, stat_val,
                      font_size=fsize, bold=True, color=ms_blue)
        _add_text_box(slide, 2.9, y + 0.05, 3.3, 0.5, desc,
                      font_size=12, color="#333333")
        y += 0.85

    # Right half — peer reference cards
    _add_text_box(slide, 7.0, 1.8, 5.5, 0.35,
                  "What Your Peers Are Seeing", font_size=14,
                  bold=True, color="#333333")
    peer_top = 2.3
    for peer_text in copilot.get("peers", []):
        _add_info_card(slide, 7.0, peer_top, 5.5, 0.9, peer_text)
        peer_top += 1.1

    # ── Slide 4: Your Data Platform in Action ────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BLANK_HEAD])
    _add_accent_bar(slide, 0, 0, 13.333, 0.04, primary)
    _set_placeholder_text(slide, 0, "Your Data Platform in Action",
                          font_size=28, bold=True)

    fabric = _get_fabric_highlights(fabric_iq, foundry_iq)

    # Left — big stat hero
    _add_text_box(slide, 0.8, 1.9, 3.5, 0.7, fabric["monthly_queries"],
                  font_size=44, bold=True, color=ms_blue)
    _add_text_box(slide, 0.8, 2.55, 3.5, 0.4, "queries per month",
                  font_size=14, color="#333333")

    workload_text = "\n".join(f"\u2022 {w}" for w in fabric["workloads"])
    _add_text_box(slide, 0.8, 3.2, 5.0, 1.2, workload_text,
                  font_size=12, color="#606060")

    # Right — Scaling for Growth
    _add_text_box(slide, 7.0, 1.9, 5.5, 0.35, "Scaling for Growth",
                  font_size=14, bold=True, color="#333333")

    growth_text = (
        f"With {fabric['utilization_pct']}% capacity utilization, your "
        f"analytics workloads are growing rapidly. We are actively "
        f"optimizing your platform performance."
    )
    _add_text_box(slide, 7.0, 2.35, 5.5, 0.7, growth_text,
                  font_size=11, color="#606060")

    # "Our Commitment" card — light-blue background
    _add_info_card(
        slide, 7.0, 3.2, 5.5, 1.1,
        f"Our Commitment\n\n{fabric['commitment']}",
        font_size=11, bg_color=(220, 235, 250),
    )

    # ── Slide 5: Opportunities Ahead ─────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_3COL_SUBHEADS])
    _add_accent_bar(slide, 0, 0, 13.333, 0.04, primary)
    _set_placeholder_text(slide, 0, "Opportunities Ahead",
                          font_size=28, bold=True)

    opps = _get_opportunities(fabric_iq, foundry_iq)
    col_map = [(16, 15), (22, 23), (24, 25)]
    for i, (sub_idx, body_idx) in enumerate(col_map):
        if i < len(opps):
            _set_placeholder_text(slide, sub_idx, opps[i]["subhead"],
                                  font_size=14, bold=True, color=ms_blue)
            _set_placeholder_text(slide, body_idx, opps[i]["body"],
                                  font_size=11)

    # ── Slide 6: Next Steps + Close ──────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_1COL_TEXT])
    _add_accent_bar(slide, 0, 0, 13.333, 0.04, primary)
    _set_placeholder_text(slide, 0, "Recommended Next Steps",
                          font_size=28, bold=True)

    # Subhead — try the placeholder; not all layouts expose it
    try:
        _set_placeholder_text(slide, 16, "How we move forward together",
                              font_size=16, color="#606060")
    except KeyError:
        pass

    steps = _get_next_steps()
    _set_placeholder_text(slide, 15, "\n\n".join(steps), font_size=13)

    # Closing bar with thank-you
    bar_top = 6.8
    _add_accent_bar(slide, 0, bar_top, 13.333, 0.5, primary)
    _add_text_box(
        slide, 0.5, bar_top + 0.05, 12.333, 0.4,
        f"Thank you for your continued partnership, {brand['display_name']}.",
        font_size=14, bold=True, color="#FFFFFF",
        alignment=PP_ALIGN.CENTER,
    )

    # -- Save --
    filename = f"presentation_{brand['short_name'].lower().replace(' ', '_')}_{date.today().isoformat()}.pptx"
    path = out_dir / filename
    prs.save(str(path))
    return str(path)
