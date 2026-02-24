"""Tests for document generation."""

import json
import os

import pytest

from src.config import MOCK_DATA_DIR


@pytest.fixture(autouse=True)
def _mock_mode():
    os.environ["USE_MOCK_DATA"] = "true"


@pytest.fixture
def coca_cola_data():
    """Load all mock data for Coca-Cola."""
    with open(MOCK_DATA_DIR / "work_iq_data.json") as f:
        work_iq = json.load(f)["coca-cola"]
    with open(MOCK_DATA_DIR / "fabric_iq_data.json") as f:
        fabric_iq = json.load(f)["coca-cola"]
    with open(MOCK_DATA_DIR / "foundry_iq_data.json") as f:
        foundry_iq = json.load(f)["coca-cola"]
    return work_iq, fabric_iq, foundry_iq


def test_generate_prep_doc(coca_cola_data):
    from src.tools.doc_generator import generate_prep_doc

    work_iq, fabric_iq, foundry_iq = coca_cola_data
    path = generate_prep_doc("Coca-Cola", work_iq, fabric_iq, foundry_iq)

    assert path.endswith(".docx")
    assert os.path.exists(path)
    assert os.path.getsize(path) > 0

    # Verify it's a valid docx
    from docx import Document
    doc = Document(path)
    full_text = "\n".join(p.text for p in doc.paragraphs)
    assert "Coca-Cola" in full_text
    assert "Relationship Context" in full_text
    assert "Business Health" in full_text


def test_generate_presentation(coca_cola_data):
    from src.tools.doc_generator import generate_presentation

    work_iq, fabric_iq, foundry_iq = coca_cola_data
    path = generate_presentation("Coca-Cola", work_iq, fabric_iq, foundry_iq)

    assert path.endswith(".pptx")
    assert os.path.exists(path)
    assert os.path.getsize(path) > 0

    # Verify it's a valid pptx with 6 slides
    from pptx import Presentation
    prs = Presentation(path)
    assert len(prs.slides) == 6


def test_generate_presentation_uses_template(coca_cola_data):
    """Verify the presentation uses Microsoft Brand Template layouts (not blank)."""
    from src.tools.doc_generator import generate_presentation
    from pptx import Presentation

    work_iq, fabric_iq, foundry_iq = coca_cola_data
    path = generate_presentation("Coca-Cola", work_iq, fabric_iq, foundry_iq)
    prs = Presentation(path)

    # Template slides have layout names from the .potx
    layout_names = [slide.slide_layout.name for slide in prs.slides]
    assert "1_Title_Gradient_Warm Gray" in layout_names
    assert "Blank_with Head" in layout_names
    assert "3-column_Text_with Subheads" in layout_names
    assert "1_1-column_Text" in layout_names


def test_prep_doc_has_all_sections(coca_cola_data):
    from src.tools.doc_generator import generate_prep_doc
    from docx import Document

    work_iq, fabric_iq, foundry_iq = coca_cola_data
    path = generate_prep_doc("Coca-Cola", work_iq, fabric_iq, foundry_iq)
    doc = Document(path)
    full_text = "\n".join(p.text for p in doc.paragraphs)

    assert "Relationship Context" in full_text
    assert "Business Health" in full_text
    assert "Recommended Discussion Topics" in full_text
    assert "Relevant Sales Plays" in full_text


def test_presentation_no_internal_data(coca_cola_data):
    """Verify no internal-only data appears in the customer-facing deck."""
    from src.tools.doc_generator import generate_presentation
    from pptx import Presentation

    work_iq, fabric_iq, foundry_iq = coca_cola_data
    path = generate_presentation("Coca-Cola", work_iq, fabric_iq, foundry_iq)
    prs = Presentation(path)

    # Collect all text from the presentation
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    full_text = "\n".join(all_text)

    # These internal items must NOT appear
    forbidden = [
        "Marcus Williams",       # account team
        "Priya Sharma",          # account team
        "James O'Brien",         # account team
        "Renewal Risk",          # internal metric
        "TKT-",                  # ticket IDs
        "Confidence:",           # internal score
        "Proposed Annual Spend", # internal revenue
        "Microsoft Internal",    # confidentiality label
        "PepsiCo",               # must be anonymized
        "Nestlé",                # must be anonymized
        "AB InBev",              # must be anonymized
    ]
    for term in forbidden:
        assert term not in full_text, (
            f"Internal data leaked: '{term}' found in presentation"
        )


def test_presentation_has_customer_branding(coca_cola_data):
    """Check that the presentation has Coca-Cola red (#FE0000) accent bars."""
    from src.tools.doc_generator import generate_presentation
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    work_iq, fabric_iq, foundry_iq = coca_cola_data
    path = generate_presentation("Coca-Cola", work_iq, fabric_iq, foundry_iq)
    prs = Presentation(path)

    # Check all slides for the Coca-Cola red accent bar
    found_red = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "fill") and shape.fill.type is not None:
                try:
                    color = shape.fill.fore_color.rgb
                    if color == RGBColor(0xFE, 0x00, 0x00):
                        found_red = True
                        break
                except Exception:
                    continue
        if found_red:
            break
    assert found_red, "Expected Coca-Cola red (#FE0000) accent bar on slides"


def test_generate_presentation_accepts_json_strings(coca_cola_data):
    """Function tool pattern: parameters can be JSON strings."""
    from src.tools.doc_generator import generate_presentation

    work_iq, fabric_iq, foundry_iq = coca_cola_data
    path = generate_presentation(
        "Coca-Cola",
        json.dumps(work_iq),
        json.dumps(fabric_iq),
        json.dumps(foundry_iq),
    )
    assert path.endswith(".pptx")
    assert os.path.exists(path)


def test_generate_prep_doc_accepts_json_strings(coca_cola_data):
    """Function tool pattern: parameters can be JSON strings."""
    from src.tools.doc_generator import generate_prep_doc

    work_iq, fabric_iq, foundry_iq = coca_cola_data
    path = generate_prep_doc(
        "Coca-Cola",
        json.dumps(work_iq),
        json.dumps(fabric_iq),
        json.dumps(foundry_iq),
    )
    assert path.endswith(".docx")
    assert os.path.exists(path)
